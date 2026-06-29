import os
import tempfile
from unittest.mock import MagicMock, patch

from lazyllm.tools.writer.data_models import (
    DocBlock,
    DocIR,
    ResourceProfile,
    WritingContext,
    WritingOutput,
    WritingTask,
)
from lazyllm.tools.writer.data_models.task import InputResource
from lazyllm.tools.writer.tools.resource_tools import WriterResourceTools, _read_resource_content
from lazyllm.tools.writer.tools.context_tools import WriterContextTools
from lazyllm.tools.writer.utils import load_artifact_json


# ---------------------------------------------------------------------------
# helpers
# ---------------------------------------------------------------------------

def _make_doc_adapter(blocks=None, plain_text="第一段\n第二段",
                      doc_id="doc-1", ref_title="飞书文档"):
    adapter = MagicMock()
    adapter.resolve_link.return_value = {
        "provider": "feishu", "object_id": doc_id,
        "object_type": "docx", "title": ref_title, "has_child": False,
    }
    adapter.read_bytes.return_value = plain_text.encode("utf-8")
    adapter.get_document_id.return_value = doc_id
    adapter.get_doc_blocks.return_value = blocks or [
        {"block_id": "b1", "block_type": "heading", "plain_text": "标题"},
        {"block_id": "b2", "block_type": "paragraph", "plain_text": "正文"},
    ]
    return adapter


def _call_target_to_doc_ir(adapter, artifact_store, protocol="feishu", **kwargs):
    with patch("lazyllm.tools.fs.client.FS._parse",
               return_value=(protocol, None, "~docx/doc-1")):
        with patch("lazyllm.tools.fs.client.FS._get_or_create_fs",
                   return_value=adapter):
            tool = WriterResourceTools(artifact_store=artifact_store)
            params = {"uri": "feishu://~docx/doc-1", "adapter": "feishu",
                      "title": "飞书文档", "doc_id": "doc-1",
                      "meta": {"space_id": "wikcn-demo"}, **kwargs}
            return tool.target_to_doc_ir(target_document=params)


def _call_write_to_document(adapter, markdown, artifact_store,
                            protocol="feishu", uri="feishu:///write-test.md", **kwargs):
    with patch("lazyllm.tools.fs.client.FS._parse",
               return_value=(protocol, None, "/write-test.md")):
        with patch("lazyllm.tools.fs.client.FS._get_or_create_fs",
                   return_value=adapter):
            tool = WriterResourceTools(artifact_store=artifact_store)
            params = {"uri": uri, "adapter": "feishu", **kwargs}
            return tool.write_to_document(markdown=markdown,
                                          target_document=params)


# ---------------------------------------------------------------------------
# context-tools tests (pre-existing)
# ---------------------------------------------------------------------------

def test_create_writing_context_tool_result():
    task = WritingTask(task_id="task-1", query="写一份产品方案", task_type="write")
    profiles = [
        ResourceProfile(
            resource_id="res-1",
            resource_role="background",
            summary="产品背景资料",
            key_facts=["支持私有化部署"],
            style_notes=["正式"],
        )
    ]
    doc_ir = DocIR(
        doc_id="doc-1",
        blocks=[
            DocBlock(block_id="b1", block_type="heading", text="方案背景", level=1),
            DocBlock(block_id="b2", block_type="paragraph", text="这里是背景内容。"),
        ],
    )

    with tempfile.TemporaryDirectory() as d:
        tool = WriterContextTools(artifact_store=d)
        result = tool.create_writing_context(
            task=task.model_dump(),
            resource_profiles=[profile.model_dump() for profile in profiles],
            doc_ir=doc_ir.model_dump(),
        )

        assert result["artifact_path"].endswith("writing_context.json")
        assert result["context_path"] == result["artifact_path"]
        assert result["metadata"]["step_name"] == "create_writing_context"
        assert result["metadata"]["artifact_key"] == "writing_context"
        assert result["metadata"]["counts"]["facts"] == 1

        context = load_artifact_json(result["context_path"], WritingContext)
        assert context.context_id == "task-1"
        assert context.doc_id == "doc-1"
        assert context.facts[0].value == "支持私有化部署"
        assert len(context.block_summaries) == 2


def test_update_writing_context_tool_result_from_paths():
    context = WritingContext(context_id="ctx-1")
    output = WritingOutput(title="最终稿", content="这是最终输出内容。")

    with tempfile.TemporaryDirectory() as d:
        context_path = os.path.join(d, "context.json")
        output_path = os.path.join(d, "output.json")
        context.save(context_path)
        output.save(output_path)

        tool = WriterContextTools(artifact_store=d)
        result = tool.update_writing_context(content_artifact=output_path, context=context_path)

        assert result["artifact_path"].endswith("writing_context.json")
        assert result["metadata"]["step_name"] == "update_writing_context"
        assert result["metadata"]["counts"]["block_summaries"] == 1

        updated = load_artifact_json(result["context_path"], WritingContext)
        assert updated.document_summary.summary == "最终稿 这是最终输出内容。"
        assert updated.block_summaries[0].summary == "最终稿 这是最终输出内容。"


# ---------------------------------------------------------------------------
# target_to_doc_ir
# ---------------------------------------------------------------------------

def test_target_to_doc_ir():
    """Normal path: produces DocIR artifact with correct fields and source."""
    adapter = _make_doc_adapter()

    with tempfile.TemporaryDirectory() as d:
        result = _call_target_to_doc_ir(adapter, d)

        assert result["artifact_path"].endswith("doc_ir.json")
        assert result["metadata"]["step_name"] == "target_to_doc_ir"
        assert result["metadata"]["counts"]["blocks"] == 2
        assert result["metadata"]["extra"]["adapter"] == "feishu"

        doc_ir = load_artifact_json(result["artifact_path"], DocIR)
        assert doc_ir.doc_id == "doc-1"
        assert doc_ir.title == "飞书文档"
        assert doc_ir.adapter == "feishu"
        assert doc_ir.plain_text == "第一段\n第二段"
        assert doc_ir.blocks[0].block_type == "heading"
        assert doc_ir.blocks[1].block_type == "paragraph"

        # source kept as-is (no back-fill)
        assert doc_ir.source.uri == "feishu://~docx/doc-1"
        assert doc_ir.source.adapter == "feishu"
        assert "source_locator" not in doc_ir.meta

        # adapter called in correct order
        adapter.resolve_link.assert_called_once()
        adapter.read_bytes.assert_called_once()
        adapter.get_document_id.assert_called_once()
        adapter.get_doc_blocks.assert_called_once()
        calls = [c[0] for c in adapter.method_calls]
        assert calls.index("resolve_link") < calls.index("read_bytes")
        assert calls.index("read_bytes") < calls.index("get_document_id")
        assert calls.index("get_document_id") < calls.index("get_doc_blocks")


def test_target_to_doc_ir_source_snapshot():
    """source keeps user value even when title comes from resolved_ref."""
    adapter = _make_doc_adapter(ref_title="Resolved Title")

    with tempfile.TemporaryDirectory() as d:
        result = _call_target_to_doc_ir(adapter, d, title="User Title")
        doc_ir = load_artifact_json(result["artifact_path"], DocIR)

        assert doc_ir.source.doc_id == "doc-1"
        assert doc_ir.title == "User Title"
        assert doc_ir.source.title == "User Title"


def test_target_to_doc_ir_title_fallback():
    """title falls back to resolved_ref when target does not supply one."""
    adapter = _make_doc_adapter(ref_title="Resolved Title")

    with tempfile.TemporaryDirectory() as d:
        result = _call_target_to_doc_ir(adapter, d, title=None)
        doc_ir = load_artifact_json(result["artifact_path"], DocIR)

    assert doc_ir.title == "Resolved Title"


def test_target_to_doc_ir_normalizes_int_block_types():
    """Feishu int block_types map to DocIR Literals with unknown→\"block\"."""
    adapter = _make_doc_adapter(blocks=[
        {"block_id": "b1", "block_type": 1,  "plain_text": ""},
        {"block_id": "b2", "block_type": 3,  "plain_text": "标题"},
        {"block_id": "b3", "block_type": 2,  "plain_text": "正文"},
        {"block_id": "b4", "block_type": 14, "plain_text": "代码"},
        {"block_id": "b5", "block_type": 31, "plain_text": ""},
        {"block_id": "b6", "block_type": 99, "plain_text": "未知"},
    ], plain_text="全文")

    with tempfile.TemporaryDirectory() as d:
        result = _call_target_to_doc_ir(adapter, d)
        doc_ir = load_artifact_json(result["artifact_path"], DocIR)

    assert [b.block_type for b in doc_ir.blocks] == [
        "document", "heading", "paragraph", "code", "table", "block",
    ]


def test_target_to_doc_ir_non_document_fs():
    """Plain FS (protocol=file) does not crash — hasattr guards work."""
    class _PlainFS:
        def read_bytes(self, *_):
            return b"plain"

    with patch("lazyllm.tools.fs.client.FS._parse",
               return_value=("file", None, "/tmp/doc.txt")):
        with patch("lazyllm.tools.fs.client.FS._get_or_create_fs",
                   return_value=_PlainFS()):
            with tempfile.TemporaryDirectory() as d:
                tool = WriterResourceTools(artifact_store=d)
                result = tool.target_to_doc_ir(target_document={
                    "uri": "/tmp/doc.txt", "adapter": "file", "title": "plain",
                })
                doc_ir = load_artifact_json(result["artifact_path"], DocIR)

    assert doc_ir.plain_text == "plain"
    assert len(doc_ir.blocks) == 0
    assert doc_ir.doc_id is None


# ---------------------------------------------------------------------------
# write_to_document
# ---------------------------------------------------------------------------

def test_write_to_document():
    """markdown is written via fs.write_file; write_result artifact saved."""
    adapter = _make_doc_adapter()

    with tempfile.TemporaryDirectory() as d:
        result = _call_write_to_document(adapter, "# Hello\n\nworld", d)

        assert result["artifact_path"].endswith("write_result.json")
        assert result["metadata"]["step_name"] == "write_to_document"
        assert result["metadata"]["extra"]["adapter"] == "feishu"
        assert result["metadata"]["extra"]["document_id"] == "doc-1"

        adapter.write_file.assert_called_once()
        args = adapter.write_file.call_args[0]
        assert "write-test" in args[0]
        assert b"Hello" in args[1]


def test_write_to_document_non_document_fs():
    """Plain FS does not crash — resolve_link guard works, doc_id is empty."""
    class _PlainWriteFS:
        def write_file(self, *_):
            pass

    with patch("lazyllm.tools.fs.client.FS._parse",
               return_value=("file", None, "/tmp/doc.md")):
        with patch("lazyllm.tools.fs.client.FS._get_or_create_fs",
                   return_value=_PlainWriteFS()):
            with tempfile.TemporaryDirectory() as d:
                tool = WriterResourceTools(artifact_store=d)
                result = tool.write_to_document(
                    markdown="# Hi",
                    target_document={"uri": "/tmp/doc.md", "adapter": "file"},
                )

    assert result["metadata"]["extra"]["document_id"] == ""


# ---------------------------------------------------------------------------
# _read_resource_content
# ---------------------------------------------------------------------------

def test_read_content_text():
    res = InputResource(resource_type="text", inline_text="本产品要求支持私有化部署")
    assert _read_resource_content(res) == "本产品要求支持私有化部署"


def test_read_content_text_empty():
    res = InputResource(resource_type="text", inline_text="")
    assert _read_resource_content(res) == ""


def test_read_content_document():
    res = InputResource(resource_type="document", uri="feishu://~docx/doc-1")

    class _FakeFS:
        def read_bytes(self, *_):
            return b"# Title\n\n## Content"

    with patch("lazyllm.tools.fs.client.FS._parse", return_value=("feishu", None, "~docx/doc-1")):
        with patch("lazyllm.tools.fs.client.FS._get_or_create_fs", return_value=_FakeFS()):
            assert _read_resource_content(res) == "# Title\n\n## Content"


def test_read_content_file_pdf():
    from pypdf import PdfWriter
    from pathlib import Path
    import tempfile

    writer = PdfWriter()
    writer.add_blank_page(width=612, height=792)
    with tempfile.TemporaryDirectory() as d:
        path = Path(d) / "test.pdf"
        with open(path, "wb") as f:
            writer.write(f)
        res = InputResource(resource_type="file", uri=str(path))
        content = _read_resource_content(res)
    assert isinstance(content, str)


def test_read_content_file_markdown():
    import tempfile
    from pathlib import Path
    with tempfile.TemporaryDirectory() as d:
        path = Path(d) / "test.md"
        path.write_text("# Hello\n\nMarkdown content.", encoding="utf-8")
        res = InputResource(resource_type="file", uri=str(path))
        content = _read_resource_content(res)
    assert "Markdown content" in content
    assert "# Hello" in content


def test_read_content_file_docx():
    from docx import Document
    import tempfile
    from pathlib import Path
    doc = Document()
    doc.add_heading("Test Heading", level=1)
    doc.add_paragraph("This is a paragraph in docx.")
    with tempfile.TemporaryDirectory() as d:
        path = Path(d) / "test.docx"
        doc.save(str(path))
        res = InputResource(resource_type="file", uri=str(path))
        content = _read_resource_content(res)
    assert "Test Heading" in content
    assert "paragraph in docx" in content


def test_read_content_file_pptx():
    import pytest
    try:
        from pptx import Presentation
    except ImportError:
        pytest.skip("python-pptx not installed")

    import tempfile
    from pathlib import Path
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "PPTX Title"
    with tempfile.TemporaryDirectory() as d:
        path = Path(d) / "test.pptx"
        prs.save(str(path))
        res = InputResource(resource_type="file", uri=str(path))
        try:
            content = _read_resource_content(res)
        except ImportError:
            pytest.skip("transformers not installed for PPTXReader")
    assert "PPTX Title" in content


def test_read_content_file_csv():
    import pandas as pd
    import tempfile
    from pathlib import Path
    df = pd.DataFrame({"name": ["张三", "李四"], "age": [30, 25]})
    with tempfile.TemporaryDirectory() as d:
        path = Path(d) / "test.csv"
        df.to_csv(path, index=False)
        res = InputResource(resource_type="file", uri=str(path))
        content = _read_resource_content(res)
    assert isinstance(content, str)
    assert "张三" in content or "李四" in content


def test_read_content_image():
    from PIL import Image
    import tempfile
    from pathlib import Path
    img = Image.new("RGB", (10, 10), color="red")
    with tempfile.TemporaryDirectory() as d:
        path = Path(d) / "test.png"
        img.save(str(path))
        res = InputResource(resource_type="image", uri=str(path))
        content = _read_resource_content(res)
    assert content == ""  # no OCR by default


def test_read_content_image_with_summary():
    res = InputResource(resource_type="image", uri="/tmp/img.png",
                        summary="红色方块图片")
    assert _read_resource_content(res) == "红色方块图片"


def test_read_content_url_fallback():
    res = InputResource(resource_type="url",
                        uri="https://example.com",
                        summary="网页摘要")
    assert _read_resource_content(res) == "网页摘要"


def test_read_content_kb_fallback():
    res = InputResource(resource_type="kb",
                        kb_id="kb-123",
                        summary="知识库检索结果摘要")
    assert _read_resource_content(res) == "知识库检索结果摘要"


def test_read_content_no_summary_fallback():
    res = InputResource(resource_type="url", uri="https://example.com")
    assert _read_resource_content(res) == ""


# ---------------------------------------------------------------------------
# profile_resources
# ---------------------------------------------------------------------------

def _make_llm_result(**overrides):
    return ResourceProfile(
        resource_id="r1",
        resource_role=overrides.pop("resource_role", "spec"),
        template_usage=overrides.pop("template_usage", "structure"),
        summary=overrides.pop("summary", "LLM generated summary"),
        key_facts=overrides.pop("key_facts", ["fact1", "fact2"]),
        style_notes=overrides.pop("style_notes", ["formal", "technical"]),
        confidence=overrides.pop("confidence", 0.9),
        extracted_constraints=overrides.pop("extracted_constraints", {"format": "markdown"}),
        extracted_outline=overrides.pop("extracted_outline", None),
        **overrides,
    )


def test_profile_resources_text_no_llm():
    task = WritingTask(query="写方案", task_type="write")
    res = InputResource(resource_type="text", inline_text="需求文档",
                        summary="用户摘要", resource_id="r1",
                        meta={"role": "spec", "template": "structure"})

    with tempfile.TemporaryDirectory() as d:
        tool = WriterResourceTools(artifact_store=d)
        result = tool.profile_resources(
            task=task.model_dump(),
            input_resources=[res.model_dump()],
        )

        assert result["metadata"]["counts"]["resource_profiles"] == 1
        with open(result["artifact_path"]) as f:
            import json
            profiles = json.load(f)["data"]

        assert profiles[0]["resource_id"] == "r1"
        assert profiles[0]["resource_role"] == "spec"
        assert profiles[0]["template_usage"] == "structure"
        assert profiles[0]["summary"] == "用户摘要"
        assert profiles[0]["key_facts"] == []
        assert profiles[0]["style_notes"] == []
        assert profiles[0]["confidence"] == 1.0
        assert profiles[0]["extracted_constraints"] == {}
        assert profiles[0]["extracted_outline"] is None


def test_profile_resources_text_with_llm():
    task = WritingTask(query="写方案", task_type="write")
    res = InputResource(resource_type="text", inline_text="需求文档",
                        resource_id="r1")

    llm_result = _make_llm_result(
        resource_role="spec", template_usage="both",
        summary="LLM summary", key_facts=["hello", "world"],
        style_notes=["casual"], confidence=0.85,
        extracted_constraints={"word_limit": "5000"},
    )

    with tempfile.TemporaryDirectory() as d:
        tool = WriterResourceTools(artifact_store=d, llm=MagicMock())
        with patch.object(tool, "_call_llm_structured", return_value=llm_result):
            result = tool.profile_resources(
                task=task.model_dump(),
                input_resources=[res.model_dump()],
            )

        with open(result["artifact_path"]) as f:
            import json
            profiles = json.load(f)["data"]

        p = profiles[0]
        assert p["resource_role"] == "spec"
        assert p["template_usage"] == "both"
        assert p["summary"] == "LLM summary"
        assert p["key_facts"] == ["hello", "world"]
        assert p["style_notes"] == ["casual"]
        assert p["confidence"] == 0.85
        assert p["extracted_constraints"] == {"word_limit": "5000"}


def test_profile_resources_llm_exception_fallback():
    task = WritingTask(query="写方案", task_type="write")
    res = InputResource(resource_type="text", inline_text="需求文档",
                        summary="fallback summary", resource_id="r1")

    with tempfile.TemporaryDirectory() as d:
        tool = WriterResourceTools(artifact_store=d, llm=MagicMock())
        with patch.object(tool, "_call_llm_structured",
                          side_effect=RuntimeError("LLM down")):
            result = tool.profile_resources(
                task=task.model_dump(),
                input_resources=[res.model_dump()],
            )

        with open(result["artifact_path"]) as f:
            import json
            profiles = json.load(f)["data"]

        p = profiles[0]
        assert p["resource_role"] == "background"
        assert p["summary"] == "fallback summary"
        assert p["key_facts"] == []
        assert p["confidence"] == 1.0


def test_profile_resources_empty_inputs():
    task = WritingTask(query="写方案", task_type="write")

    with tempfile.TemporaryDirectory() as d:
        tool = WriterResourceTools(artifact_store=d)
        result = tool.profile_resources(
            task=task.model_dump(),
            input_resources=[],
        )

        assert result["metadata"]["counts"]["resource_profiles"] == 0
        with open(result["artifact_path"]) as f:
            import json
            profiles = json.load(f)["data"]
        assert profiles == []


def test_profile_resources_multi():
    task = WritingTask(query="写方案", task_type="write")
    resources = [
        InputResource(resource_type="text", inline_text="需求文档",
                      resource_id="r1"),
        InputResource(resource_type="text", inline_text="背景材料",
                      resource_id="r2"),
        InputResource(resource_type="text", inline_text="范文示例",
                      resource_id="r3"),
    ]

    with tempfile.TemporaryDirectory() as d:
        tool = WriterResourceTools(artifact_store=d)
        result = tool.profile_resources(
            task=task.model_dump(),
            input_resources=[r.model_dump() for r in resources],
        )

        assert result["metadata"]["counts"]["resource_profiles"] == 3
        with open(result["artifact_path"]) as f:
            import json
            profiles = json.load(f)["data"]
        assert [p["resource_id"] for p in profiles] == ["r1", "r2", "r3"]


def test_profile_resources_meta_overrides():
    task = WritingTask(query="写方案", task_type="write")
    res = InputResource(resource_type="text", inline_text="内容",
                        resource_id="r1",
                        meta={"role": "example", "template": "style"})

    with tempfile.TemporaryDirectory() as d:
        tool = WriterResourceTools(artifact_store=d)
        result = tool.profile_resources(
            task=task.model_dump(),
            input_resources=[res.model_dump()],
        )

        with open(result["artifact_path"]) as f:
            import json
            profiles = json.load(f)["data"]

        p = profiles[0]
        assert p["resource_role"] == "example"
        assert p["template_usage"] == "style"


def test_profile_resources_content_fallback():
    task = WritingTask(query="写方案", task_type="write")
    res = InputResource(resource_type="url", uri="https://x.com",
                        resource_id="r1")

    with tempfile.TemporaryDirectory() as d:
        tool = WriterResourceTools(artifact_store=d)
        result = tool.profile_resources(
            task=task.model_dump(),
            input_resources=[res.model_dump()],
        )

        with open(result["artifact_path"]) as f:
            import json
            profiles = json.load(f)["data"]

        p = profiles[0]
        assert p["summary"] == ""


def test_profile_resources_artifact_structure():
    task = WritingTask(query="写方案", task_type="write",
                       constraints={"word_limit": "5000"})
    res = InputResource(resource_type="text", inline_text="需求文档")

    with tempfile.TemporaryDirectory() as d:
        tool = WriterResourceTools(artifact_store=d)
        result = tool.profile_resources(
            task=task.model_dump(),
            input_resources=[res.model_dump()],
        )

        assert result["artifact_path"].endswith("resource_profiles.json")
        assert result["metadata"]["step_name"] == "profile_resources"
        assert result["metadata"]["artifact_key"] == "resource_profiles"
        assert result["metadata"]["status"] == "success"

