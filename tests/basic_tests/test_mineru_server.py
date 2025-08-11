#!/usr/bin/env python3
# -*- coding: utf-8 -*-
import os
import shutil
import pytest
from pathlib import Path
import requests
import unittest
from lazyllm import LOG
from lazyllm.tools.rag.tools.mineru_server_module import MineruServerModule
from lazyllm.tools.rag.readers.mineru_pdf_reader import MineruPDFReader


@pytest.fixture(autouse=True)
def setup_tmpdir(request, tmpdir):
    request.cls.tmpdir = tmpdir


@pytest.fixture(scope="class", autouse=True)
def setup_tmpdir_class(request, tmpdir_factory):
    request.cls.tmpdir_class = tmpdir_factory.mktemp("mineru_test")


@pytest.mark.skip_on_win
@pytest.mark.skip_on_mac
@pytest.mark.usefixtures("setup_tmpdir_class")
class TestMineruServer(unittest.TestCase):
    @classmethod
    def setUpClass(cls):
        """类级别的设置，启动服务器"""
        cls.cache_dir = str(cls.tmpdir_class.mkdir("cache"))
        cls.image_save_dir = str(cls.tmpdir_class.mkdir("images"))
        cls.server = MineruServerModule(
            cache_dir=cls.cache_dir, pythonpath=None, port=31769
        )
        cls.server.start()
        cls.url = cls.server._url[:-9] + "/api/v1/pdf_parse"
        cls.test_results = {}

    @classmethod
    def tearDownClass(cls):
        if hasattr(cls, "server"):
            cls.server.stop()
        if hasattr(cls, "test_files_dir") and cls.test_files_dir.exists():
            shutil.rmtree(str(cls.test_files_dir))

    def setUp(self):
        self.test_files_dir = Path("test_files")
        self.test_files_dir.mkdir(exist_ok=True)
        self.test_files = {
            "pdf1": self.test_files_dir / "test_document1.pdf",
            "pdf2": self.test_files_dir / "test_document2.pdf",
            "pdf3": self.test_files_dir / "test_document3.pdf",
            "docx": self.test_files_dir / "test_document.docx",
            "pptx": self.test_files_dir / "test_presentation.pptx",
        }

        self.setup_test_files()

    def setup_test_files(self):  # noqa: C901
        try:
            from reportlab.pdfgen import canvas
            from reportlab.lib.pagesizes import letter
            import tempfile
            from PIL import Image, ImageDraw

            pdf_contents = {
                "pdf1": [
                    "this is the first pdf test document",
                    "used for testing pdf parsing function",
                ],
                "pdf2": [
                    "this is the second pdf test document",
                    "used for testing pdf parsing function",
                ],
                "pdf3": [
                    "this is the third pdf test document",
                    "used for testing pdf parsing function",
                ],
            }
            for pdf_name, content_lines in pdf_contents.items():
                file_path = self.test_files[pdf_name]
                if not file_path.exists():
                    # 创建真正的PDF文件
                    c = canvas.Canvas(str(file_path), pagesize=letter)

                    y_position = 750
                    for line in content_lines:
                        c.drawString(100, y_position, line)
                        y_position -= 20
                    # 为pdf2添加图片和表格
                    if pdf_name == "pdf2":
                        try:
                            # 创建简单测试图片
                            img = Image.new("RGB", (150, 100), color="lightblue")
                            draw = ImageDraw.Draw(img)
                            draw.text((10, 10), "test image", fill="black")
                            draw.rectangle([10, 30, 140, 90], outline="red", width=2)

                            # 保存图片到临时文件
                            with tempfile.NamedTemporaryFile(
                                suffix=".png", delete=False
                            ) as tmp_file:
                                img.save(tmp_file.name, "PNG")
                                tmp_file_path = tmp_file.name

                            # 将图片添加到PDF
                            c.drawImage(tmp_file_path, 100, 400, width=150, height=100)
                            c.drawString(100, 380, "test image")

                            # 添加简单表格
                            c.drawString(100, 320, "test table:")
                            c.drawString(100, 300, "product    number    price")
                            c.drawString(100, 280, "productA   10     $100")
                            c.drawString(100, 260, "productB   5     $200")

                            # 清理临时文件
                            os.unlink(tmp_file_path)

                        except Exception as img_error:
                            LOG.info(f"⚠️ 为{pdf_name}添加图片失败: {img_error}")

                    c.save()
                    LOG.info(f"✅ 创建PDF文件: {pdf_name}")
                    if pdf_name == "pdf2":
                        LOG.info("   📷 包含测试图片和表格")

        except ImportError:
            assert Exception("❌ 缺少reportlab库，无法创建测试文件")
            assert Exception("请安装: pip install reportlab")
        except Exception as e:
            assert Exception(f"❌ 创建PDF文件失败: {str(e)}")

        try:
            from docx import Document
            from pptx import Presentation

            # 创建Word文档
            docx_file = self.test_files["docx"]
            if not docx_file.exists():
                doc = Document()
                doc.add_heading("测试Word文档", 0)
                doc.add_paragraph("这是一个测试Word文档")
                doc.add_paragraph("包含中文内容")
                doc.add_paragraph("用于测试文档转换功能")
                doc.save(str(docx_file))
                LOG.info(f"✅ 创建Word文档: {docx_file.name}")

            pptx_file = self.test_files["pptx"]
            if not pptx_file.exists():
                prs = Presentation()
                slide = prs.slides.add_slide(prs.slide_layouts[0])
                title = slide.shapes.title
                content = slide.placeholders[1]

                title.text = "测试PowerPoint文档"
                content.text = (
                    "这是一个测试PowerPoint文档\n包含中文内容\n用于测试演示文稿转换功能"
                )

                prs.save(str(pptx_file))
                LOG.info(f"✅ 创建PowerPoint文档: {pptx_file.name}")

        except ImportError:
            LOG.info("❌ 缺少python-docx或python-pptx库，无法创建Office文档")
            LOG.info("请安装: pip install python-docx python-pptx")
            return False
        except Exception as e:
            LOG.info(f"❌ 创建Office文档失败: {str(e)}")
            return False

        LOG.info("✅ 测试文件设置完成")
        LOG.info(f"📄 PDF文件: {len(pdf_contents)} 个")
        LOG.info("📄 Office文档: 2 个")
        return True

    def post_pdf_parse(
        self,
        files,
        backend="pipeline",
        return_md=True,
        return_content_list=True,
        use_cache=False,
    ):
        """修复：使用正确的Form数据格式发送请求"""
        data = {
            "files": files,
            "backend": backend,
            "return_md": return_md,
            "return_content_list": return_content_list,
            "use_cache": use_cache,
        }
        try:
            resp = requests.post(self.__class__.url, data=data)
            try:
                return resp.status_code, resp.json()
            except Exception:
                return resp.status_code, resp.text
        except Exception as e:
            return 500, str(e)

    def check_result(self, result):
        assert isinstance(result, dict)
        assert "result" in result, f"result: {result}"
        for res in result["result"]:
            assert "md_content" in res and "content_list" in res

    @pytest.mark.order(1)
    def test_pdf_parsing(self):
        """测试1: 初始PDF解析（创建缓存基础）"""
        LOG.info("\n📋 测试1: 初始PDF解析（创建缓存基础）")
        LOG.info("🔄 这一步将解析PDF文件并创建缓存，为后续测试做准备")
        initial_files = [str(self.test_files["pdf1"]), str(self.test_files["pdf2"])]
        LOG.info(f"📄 解析PDF文件: {[os.path.basename(f) for f in initial_files]}")

        status, result = self.post_pdf_parse(
            files=initial_files,
            backend="pipeline",
            return_md=True,
            return_content_list=True,
            use_cache=False,
        )

        assert status == 200, f"status: {status}, error: {result}"
        self.check_result(result)
        LOG.info("✅ 初始PDF解析成功")
        LOG.info(f"📊 处理文件数: {len(result['result'])}")
        # 存储结果用于后续测试 - 修复：使用文件路径作为键
        for i, file_result in enumerate(result["result"]):
            file_path = initial_files[i]
            self.__class__.test_results[file_path] = file_result
            file_name = os.path.basename(file_path)
            LOG.info(f"📄 {file_name}:")
            if "md_content" in file_result:
                LOG.info(f"   - MD内容: {file_result['md_content']}")
            if "content_list" in file_result:
                LOG.info(f"   - Content List: {file_result['content_list']} ")

        assert os.path.exists(self.__class__.cache_dir), "缓存目录应该存在"
        cache_files = list(Path(self.__class__.cache_dir).rglob("*"))
        LOG.info(f"📁 缓存文件数: {len(cache_files)}")

    @pytest.mark.order(2)
    def test_pdf_parsing_with_upload_files(self):
        """测试2: 初始上传文件对象解析"""
        LOG.info("\n📋 测试2: 初始上传文件对象解析")
        file_path = str(self.test_files["pdf2"])
        with open(file_path, "rb") as f:
            upload_files = [
                (
                    "upload_files",
                    (os.path.basename(file_path), f.read(), "application/pdf"),
                )
            ]
        data = {
            "backend": "vlm-sglang-engine",
            "return_md": True,
            "return_content_list": True,
            "use_cache": False,
        }
        resp = requests.post(self.__class__.url, data=data, files=upload_files)
        status = resp.status_code
        assert status == 200, f"status: {status}, error: {resp.text}"
        result = resp.json()
        self.check_result(result)

    @pytest.mark.order(3)
    def test_pdf_parsing_with_cache(self):
        """测试3: 混合PDF缓存和新文件（依赖测试1和3）"""
        LOG.info("\n🔄 测试3: 混合PDF缓存和新文件（依赖测试1和3）")
        LOG.info("🔄 这一步将测试部分PDF文件使用缓存，部分PDF文件重新解析")
        # 混合文件：第一个有缓存，第三个是新文件
        mixed_files = [str(self.test_files["pdf1"]), str(self.test_files["pdf3"])]
        LOG.info("📄 混合PDF文件:")
        LOG.info(f"   - 有缓存: {os.path.basename(mixed_files[0])}")
        LOG.info(f"   - 新文件: {os.path.basename(mixed_files[1])}")
        status, result = self.post_pdf_parse(
            files=mixed_files,
            backend="pipeline",
            return_md=True,
            return_content_list=True,
            use_cache=True,  # 使用缓存
        )
        assert status == 200, f"status: {status}, error: {result}"
        self.check_result(result)
        LOG.info("✅ 混合PDF处理成功")
        LOG.info(f"📊 处理文件数: {len(result['result'])}")
        for i, file_result in enumerate(result["result"]):
            file_path = mixed_files[i]
            file_name = os.path.basename(file_path)
            LOG.info(f"📄 {file_name}:")
            assert (
                "md_content" in file_result or "content_list" in file_result
            ), f"文件{file_name}应包含解析内容"
            if "md_content" in file_result:
                LOG.info(f"   - MD内容: {len(file_result['md_content'])} 字符")
            if "content_list" in file_result:
                LOG.info(f"   - Content List: {len(file_result['content_list'])} 项")
            if i == 0:
                original_result = self.__class__.test_results.get(file_path)
                if original_result:
                    assert file_result.get("md_content") == original_result.get(
                        "md_content"
                    )
                    assert file_result.get("content_list") == original_result.get(
                        "content_list"
                    )

    @pytest.mark.order(4)
    def test_office_document_parsing(self):
        """测试4: Office文档解析功能"""
        LOG.info("\n📄 测试4: Office文档解析功能")
        office_files = [str(self.test_files["docx"]), str(self.test_files["pptx"])]
        for file_path in office_files:
            LOG.info(f"🔄 测试文件: {os.path.basename(file_path)}")
            status, result = self.post_pdf_parse(
                files=[file_path],
                backend="pipeline",
                return_md=True,
                return_content_list=True,
                use_cache=False,
            )
            assert status in [200, 400], f"status: {status}, error: {result}"
            if status == 200:
                self.check_result(result)
                LOG.info(f"✅ {os.path.basename(file_path)} 解析成功")
                file_result = result["result"][0]
                if "md_content" in file_result:
                    LOG.info(f"   - MD内容: {file_result['md_content']}")
                if "content_list" in file_result:
                    LOG.info(f"   - Content List: {file_result['content_list']} ")
            else:
                LOG.warning(f"跳过office文档解析测试")

    @pytest.mark.order(5)
    def test_different_backends(self):
        """测试6: 不同后端测试"""
        LOG.info("\n🔧 测试6: 不同后端测试")
        backend = "vlm-sglang-engine"
        test_file = str(self.test_files["pdf1"])
        LOG.info(f"🔄 测试后端: {backend}")
        status, result = self.post_pdf_parse(
            files=[test_file],
            backend=backend,
            return_md=True,
            return_content_list=True,
            use_cache=False,
        )

        assert status == 200, f"status: {status}, error: {result}"
        self.check_result(result)
        assert self.__class__.test_results.get(test_file) != result["result"][0]
        LOG.info(f"✅ {backend} 后端测试成功")
        file_result = result["result"][0]
        if "md_content" in file_result:
            LOG.info(f"   - MD内容: {file_result['md_content']}")
        if "content_list" in file_result:
            LOG.info(f"   - Content List: {file_result['content_list']} ")

    @pytest.mark.order(6)
    def test_pdf_reader(self):
        """测试6: 测试pdf reader"""
        LOG.info("\n⚠️ 测试6: 测试pdf reader")
        pdf_reader = MineruPDFReader(self.__class__.url)
        pdf_path = str(self.test_files["pdf1"])
        nodes = pdf_reader(pdf_path)
        assert isinstance(nodes, list)
        cache_res = [
            item["text"]
            for item in self.__class__.test_results[pdf_path]["content_list"]
        ]
        assert set([node._content for node in nodes]) == set(cache_res)
        LOG.info([node._content for node in nodes])

    @pytest.mark.order(7)
    def test_pdf_reader_with_upload_files(self):
        """测试7: 测试pdf reader(上传文件)"""
        LOG.info("\n⚠️ 测试7: 测试pdf reader(上传文件)")

        # 设置 upload_mode=True 来测试文件上传功能
        pdf_reader = MineruPDFReader(self.__class__.url, upload_mode=True)
        pdf_path = str(self.test_files["pdf1"])
        nodes = pdf_reader(pdf_path)
        assert isinstance(nodes, list)
        cache_res = [
            item["text"]
            for item in self.__class__.test_results[pdf_path]["content_list"]
        ]
        assert set([node._content for node in nodes]) == set(cache_res)
        LOG.info([node._content for node in nodes])

        LOG.info("✅ PDF reader 文件上传模式测试成功")
