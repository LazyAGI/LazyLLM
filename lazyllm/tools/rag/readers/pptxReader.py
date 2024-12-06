import os
import tempfile
from fsspec import AbstractFileSystem
from pathlib import Path
from typing import Optional, Dict, List

from .readerBase import LazyLLMReaderBase, infer_torch_device
from ..doc_node import DocNode

class PPTXReader(LazyLLMReaderBase):
    def __init__(self, return_trace: bool = True) -> None:
        try:
            import torch  # noqa
            from PIL import Image # noqa
            from pptx import Presentation  # noqa
            from transformers import (AutoTokenizer, VisionEncoderDecoderModel, ViTFeatureExtractor,)
        except ImportError:
            raise ImportError("Please install extra dependencies that are required for the "
                              "PPTXReader: `pip install torch transformers python-pptx Pillow`")

        super().__init__(return_trace=return_trace)
        model = VisionEncoderDecoderModel.from_pretrained("nlpconnect/vit-gpt2-image-captioning")
        feature_extractor = ViTFeatureExtractor.from_pretrained("nlpconnect/vit-gpt2-image-captioning")
        tokenizer = AutoTokenizer.from_pretrained("nlpconnect/vit-gpt2-image-captioning")

        self._parser_config = {"feature_extractor": feature_extractor, "model": model, "tokenizer": tokenizer}

    def _caption_image(self, tmp_image_file: str) -> str:
        from PIL import Image

        model = self._parser_config['model']
        feature_extractor = self._parser_config['feature_extractor']
        tokenizer = self._parser_config['tokenizer']

        device = infer_torch_device()
        model.to(device)

        max_length = 16
        num_beams = 4
        gen_kwargs = {"max_length": max_length, "num_beams": num_beams}

        i_image = Image.open(tmp_image_file)
        if i_image.mode != "RGB": i_image = i_image.convert(mode="RGB")

        pixel_values = feature_extractor(images=[i_image], return_tensors="pt").pixel_values
        pixel_values = pixel_values.to(device)

        output_ids = model.generate(pixel_values, **gen_kwargs)

        preds = tokenizer.batch_decode(output_ids, skip_special_tokens=True)
        return preds[0].strip()

    def _load_data(self, file: Path, extra_info: Optional[Dict] = None,
                   fs: Optional[AbstractFileSystem] = None) -> List[DocNode]:
        from pptx import Presentation

        if not isinstance(file, Path): file = Path(file)

        if fs:
            with fs.open(file) as f:
                presentation = Presentation(f)
        else:
            presentation = Presentation(file)

        result = ""
        for i, slide in enumerate(presentation.slides):
            result += f"\n\nSlide #{i}: \n"
            for shape in slide.shapes:
                if hasattr(shape, "image"):
                    image = shape.image
                    image_bytes = image.blob
                    f = tempfile.NamedTemporaryFile("wb", delete=False)
                    try:
                        f.write(image_bytes)
                        f.close()
                        result += f"\n Image: {self._caption_image(f.name)}\n\n"
                    finally:
                        os.unlink(f.name)

                if hasattr(shape, "text"): result += f"{shape.text}\n"
        return [DocNode(text=result, global_metadata=extra_info)]
